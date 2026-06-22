"""ArtifactGenerator: produces business deliverables from workflow results.

Supports: Report (markdown), Excel (.xlsx), Chart (Echarts JSON), CSV.
"""

from __future__ import annotations

import json
import os
import io
import csv
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import Optional

from backend.workflow.schemas import ArtifactType, WorkflowArtifactRef


BASE_DIR = Path(__file__).resolve().parent.parent.parent
ARTIFACT_DIR = BASE_DIR / "data" / "artifacts"


class ArtifactGenerator:
    """Generates business deliverables from workflow step results."""

    def _get_model(self):
        from backend.agent.model_router import get_model_for_agent
        return get_model_for_agent("supervisor")

    async def generate_report(
        self,
        title: str,
        step_results: dict[str, dict],
        user_context: dict | None = None,
    ) -> WorkflowArtifactRef:
        """Generate a markdown report summarizing all step results."""
        from langchain_core.messages import SystemMessage, HumanMessage

        results_text = json.dumps(step_results, ensure_ascii=False, indent=2)

        prompt = f"""Generate a professional business report in markdown based on these workflow results.

Title: {title}

Results:
{results_text[:8000]}

Structure the report with:
1. Executive Summary
2. Key Findings
3. Detailed Analysis (one section per step)
4. Recommendations
5. Appendix (raw data summary)

Use professional formatting: headings (##), bullet points, tables where appropriate."""

        model = self._get_model()
        response = await model.ainvoke([
            SystemMessage(content="You are a professional business report writer."),
            HumanMessage(content=prompt),
        ])

        content = response.content if hasattr(response, "content") else str(response)

        return WorkflowArtifactRef(
            title=title,
            artifact_type=ArtifactType.REPORT,
            mime_type="text/markdown",
            content=content,
        )

    async def generate_excel(
        self,
        title: str,
        data: list[dict],
        user_context: dict | None = None,
    ) -> WorkflowArtifactRef:
        """Generate an Excel file from structured data."""
        import pandas as pd

        os.makedirs(ARTIFACT_DIR, exist_ok=True)

        df = pd.DataFrame(data)
        filename = f"excel_{datetime.now(timezone.utc).strftime('%Y%m%d_%H%M%S')}.xlsx"
        file_path = ARTIFACT_DIR / filename

        df.to_excel(str(file_path), index=False, engine="openpyxl")

        return WorkflowArtifactRef(
            title=title,
            artifact_type=ArtifactType.EXCEL,
            mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            url=str(file_path),
        )

    async def generate_chart(
        self,
        title: str,
        data: dict,
        chart_type: str = "bar",
        user_context: dict | None = None,
    ) -> WorkflowArtifactRef:
        """Generate an Echarts JSON configuration for frontend rendering."""
        from backend.agent.chart_generator import generate_echarts_config, format_chart_markdown

        echarts_config = generate_echarts_config(data, chart_type)
        content = format_chart_markdown(echarts_config, chart_type)

        return WorkflowArtifactRef(
            title=title,
            artifact_type=ArtifactType.CHART,
            mime_type="application/json+echarts",
        )

    async def generate_csv(
        self,
        title: str,
        data: list[dict],
        user_context: dict | None = None,
    ) -> WorkflowArtifactRef:
        """Generate a CSV file from structured data."""
        os.makedirs(ARTIFACT_DIR, exist_ok=True)

        output = io.StringIO()
        if data:
            writer = csv.DictWriter(output, fieldnames=data[0].keys())
            writer.writeheader()
            writer.writerows(data)

        filename = f"csv_{datetime.now(timezone.utc).strftime('%Y%m%d_%H%M%S')}.csv"
        file_path = ARTIFACT_DIR / filename
        file_path.write_text(output.getvalue(), encoding="utf-8")

        return WorkflowArtifactRef(
            title=title,
            artifact_type=ArtifactType.CSV,
            mime_type="text/csv",
            url=str(file_path),
        )

    async def generate_pdf(self, title: str, content: str) -> WorkflowArtifactRef:
        """Generate a PDF from markdown content using reportlab."""
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        except ImportError:
            return WorkflowArtifactRef(
                artifact_type=ArtifactType.PDF,
                title=title,
                mime_type="text/plain",
                content="PDF generation requires reportlab. Install: pip install reportlab",
            )

        file_name = f"report_{uuid.uuid4().hex[:8]}.pdf"
        file_path = ARTIFACT_DIR / file_name
        os.makedirs(ARTIFACT_DIR, exist_ok=True)

        doc = SimpleDocTemplate(str(file_path), pagesize=A4,
                              rightMargin=72, leftMargin=72,
                              topMargin=72, bottomMargin=72)
        styles = getSampleStyleSheet()

        story = []
        for line in content.split("\n"):
            stripped = line.strip()
            if not stripped:
                story.append(Spacer(1, 6))
            elif stripped.startswith("# "):
                story.append(Paragraph(stripped[2:], styles["Title"]))
            elif stripped.startswith("## "):
                story.append(Paragraph(stripped[3:], styles["Heading2"]))
            elif stripped.startswith("### "):
                story.append(Paragraph(stripped[4:], styles["Heading3"]))
            elif stripped.startswith("- "):
                story.append(Paragraph(f"• {stripped[2:]}", styles["BodyText"]))
            elif stripped.startswith("> "):
                story.append(Paragraph(stripped, styles["Italic"]))
            else:
                story.append(Paragraph(stripped, styles["BodyText"]))

        doc.build(story)

        return WorkflowArtifactRef(
            step_id="report",
            artifact_type=ArtifactType.PDF,
            title=title,
            mime_type="application/pdf",
            url=str(file_path),
            content="",
        )

    async def generate_pptx(self, title: str, content: str) -> WorkflowArtifactRef:
        """Generate a PPTX presentation from markdown content using python-pptx."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            return WorkflowArtifactRef(
                artifact_type=ArtifactType.REPORT,
                title=title,
                mime_type="text/plain",
                content="PPTX generation requires python-pptx. Install: pip install python-pptx",
            )

        file_name = f"slides_{uuid.uuid4().hex[:8]}.pptx"
        file_path = ARTIFACT_DIR / file_name
        os.makedirs(ARTIFACT_DIR, exist_ok=True)

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = title

        # Content slides: split by ## headings
        sections = content.split("\n## ")
        for section in sections:
            lines = section.strip().split("\n")
            heading = lines[0].replace("# ", "").strip()
            body = "\n".join(lines[1:])

            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = heading[:100]
            if body.strip():
                slide.shapes.placeholders[1].text = body[:500]

        prs.save(str(file_path))

        return WorkflowArtifactRef(
            step_id="report",
            artifact_type=ArtifactType.REPORT,
            title=title,
            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            url=str(file_path),
            content="",
        )


_artifact_generator: Optional[ArtifactGenerator] = None


def get_artifact_generator() -> ArtifactGenerator:
    global _artifact_generator
    if _artifact_generator is None:
        _artifact_generator = ArtifactGenerator()
    return _artifact_generator
